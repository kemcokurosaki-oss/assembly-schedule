"""
組立工程表 閲覧者向け操作ガイド PowerPoint 作成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# カラーパレット（明るめ・彩度控えめトーン）
C_DARK    = RGBColor(0x3A, 0x6B, 0xAD)   # スチールブルー（タイトル・見出し背景）
C_BLUE    = RGBColor(0x5A, 0x89, 0xC4)   # ミディアムブルー（アクセント）
C_TEAL    = RGBColor(0x22, 0xA8, 0xC8)   # ソフトシアン（サブアクセント）
C_GOLD    = RGBColor(0xF5, 0x9E, 0x0B)   # アンバー（強調）
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)   # 白
C_LIGHT   = RGBColor(0xF0, 0xF6, 0xFF)   # 明るい水色系背景
C_GRAY    = RGBColor(0x47, 0x55, 0x69)   # ダークグレー（本文）
C_DARK2   = RGBColor(0x2E, 0x5F, 0x9E)   # スレートブルー
C_GREEN   = RGBColor(0x05, 0x96, 0x69)   # エメラルドグリーン（閲覧可）
C_RED     = RGBColor(0xDC, 0x26, 0x26)   # 赤（要ログインアイコン）
C_ORANGE  = RGBColor(0xD9, 0x77, 0x06)   # オレンジ（警告）

FONT_TITLE = "メイリオ"
FONT_BODY  = "メイリオ"

W = Inches(13.33)  # LAYOUT_WIDE
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def add_bg_rect(slide, x, y, w, h, color, transparency=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency is not None:
        shape.fill.fore_color.theme_color = None
        # transparency via alpha (0=opaque, 100=transparent)
        from pptx.util import Pt
        shape.fill.transparency = transparency / 100.0
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT,
             font=FONT_BODY, wrap=True, italic=False, valign=None):
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign:
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox


def add_multiline_text(slide, lines, x, y, w, h,
                       font_size=14, bold=False, color=None,
                       align=PP_ALIGN.LEFT, font=FONT_BODY,
                       line_bold=None, line_color=None, line_size=None):
    """lines: list of str or (str, bold, color, size) tuples"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line, tuple):
            txt, lb, lc, ls = (list(line) + [None, None, None, None])[:4]
        else:
            txt, lb, lc, ls = line, None, None, None

        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.name = font
        run.font.size  = Pt(ls if ls else (line_size if line_size else font_size))
        run.font.bold  = lb if lb is not None else bold
        run.font.color.rgb = lc if lc else (color if color else C_WHITE)
    return txBox


def add_badge(slide, text, x, y, w, h, bg_color, text_color=None, font_size=12, bold=True):
    rect = add_bg_rect(slide, x, y, w, h, bg_color)
    tc = text_color or C_WHITE
    add_text(slide, text, x, y, w, h,
             font_size=font_size, bold=bold, color=tc,
             align=PP_ALIGN.CENTER, font=FONT_BODY)
    return rect


def add_card(slide, x, y, w, h, bg=C_WHITE, border_color=None):
    from pptx.util import Pt
    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    if border_color:
        rect.line.color.rgb = border_color
        rect.line.width = Pt(0.75)
    else:
        rect.line.fill.background()
    return rect


# ─────────────────────────────────────────────────────────
# Slide 1: タイトル
# ─────────────────────────────────────────────────────────
def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # 全面ダーク背景
    add_bg_rect(slide, 0, 0, 13.33, 7.5, C_DARK)

    # 左ゴールドアクセントバー
    add_bg_rect(slide, 0, 0, 0.18, 7.5, C_GOLD)

    # 右装飾円
    for ci, (cx, cy, cr, alpha) in enumerate([
        (11.5, -0.5, 3.5, 0.55),
        (12.8, 2.0,  2.2, 0.60),
        (10.2, 5.5,  2.8, 0.65),
    ]):
        sh = slide.shapes.add_shape(9, Inches(cx), Inches(cy), Inches(cr), Inches(cr))
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGBColor(0x93, 0xC5, 0xFD)
        sh.fill.transparency   = alpha
        sh.line.fill.background()

    # メインタイトル
    add_text(slide, "組立工程表", 0.5, 1.6, 9.0, 1.2,
             font_size=52, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "閲覧者向け 操作ガイド", 0.5, 2.85, 9.0, 0.8,
             font_size=28, bold=False, color=C_GOLD, align=PP_ALIGN.LEFT)

    # サブテキスト
    add_text(slide, "ログインなしで使えるすべての機能を解説します",
             0.5, 4.0, 9.0, 0.6,
             font_size=16, color=RGBColor(0xBF, 0xDB, 0xFF), align=PP_ALIGN.LEFT)

    # バージョン・日付バッジ
    add_badge(slide, "閲覧者専用", 0.5, 6.5, 2.0, 0.55, C_TEAL, C_WHITE, font_size=13)

    return slide


# ─────────────────────────────────────────────────────────
# Slide 2: 画面全体の構成
# ─────────────────────────────────────────────────────────
def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)

    # ヘッダー帯
    add_bg_rect(slide, 0, 0, 13.33, 1.1, C_DARK)
    add_bg_rect(slide, 0, 0, 0.18, 1.1, C_GOLD)
    add_text(slide, "画面全体の構成", 0.4, 0.15, 10, 0.8,
             font_size=28, bold=True, color=C_WHITE)
    add_text(slide, "1 / 6", 12.0, 0.25, 1.0, 0.6,
             font_size=14, color=RGBColor(0xBA, 0xD4, 0xFF), align=PP_ALIGN.RIGHT)

    # ヘッダーエリア（ツールバー）
    add_card(slide, 0.4, 1.3, 12.5, 0.85, bg=C_DARK2)
    add_text(slide, "① ヘッダー（操作ボタン・フィルター）",
             0.6, 1.42, 10, 0.6,
             font_size=14, bold=True, color=C_WHITE)

    # グリッド
    add_card(slide, 0.4, 2.25, 3.8, 4.1, bg=C_WHITE,
             border_color=C_TEAL)
    add_bg_rect(slide, 0.4, 2.25, 3.8, 0.5, C_TEAL)
    add_text(slide, "② グリッド（左側）",
             0.5, 2.28, 3.6, 0.42,
             font_size=13, bold=True, color=C_WHITE)
    add_multiline_text(slide,
        ["工事番号・機械・ユニット",
         "タスク名・担当者",
         "組立場所・開始日・終了日",
         "日数 などの情報一覧"],
        0.55, 2.85, 3.6, 1.4,
        font_size=13, color=C_GRAY, line_color=C_GRAY)
    add_multiline_text(slide,
        [("クリック", True, C_BLUE, 13),
         ("  → タスク開始日へ自動スクロール", False, C_GRAY, 12)],
        0.55, 4.45, 3.6, 0.45,
        font_size=12, color=C_GRAY)
    add_multiline_text(slide,
        [("ダブルクリック", True, C_BLUE, 13),
         ("  → 詳細情報の確認", False, C_GRAY, 12)],
        0.55, 5.0, 3.6, 0.45,
        font_size=12, color=C_GRAY)

    # タイムライン
    add_card(slide, 4.35, 2.25, 8.55, 4.1, bg=C_WHITE,
             border_color=C_BLUE)
    add_bg_rect(slide, 4.35, 2.25, 8.55, 0.5, C_BLUE)
    add_text(slide, "③ タイムライン（右側）",
             4.45, 2.28, 8.2, 0.42,
             font_size=13, bold=True, color=C_WHITE)
    add_multiline_text(slide,
        ["日付ヘッダー（年/月/日/曜日）",
         "タスクバー（担当者ごとの色分け）",
         "今日の縦線（赤い縦線）",
         "土日・祝日は列が色分け表示"],
        4.5, 2.85, 8.0, 1.4,
        font_size=13, color=C_GRAY, line_color=C_GRAY)
    add_multiline_text(slide,
        [("左右スクロール", True, C_BLUE, 13),
         ("  → 期間を移動して確認", False, C_GRAY, 12)],
        4.5, 4.45, 8.0, 0.45,
        font_size=12, color=C_GRAY)
    add_multiline_text(slide,
        [("ダブルクリック", True, C_BLUE, 13),
         ("  → タスク詳細の確認（編集不可）", False, C_GRAY, 12)],
        4.5, 5.0, 8.0, 0.45,
        font_size=12, color=C_GRAY)

    # リサイザー矢印
    add_bg_rect(slide, 4.22, 2.25, 0.12, 4.1, RGBColor(0xCC, 0xDD, 0xEE))
    add_text(slide, "↔", 4.15, 3.8, 0.28, 0.5,
             font_size=16, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)

    # フッター注記
    add_text(slide, "※ 閲覧者はガントバーのドラッグ・編集操作はできません（ログイン必要）",
             0.4, 6.5, 12.5, 0.55,
             font_size=12, color=C_ORANGE, italic=True)

    return slide


# ─────────────────────────────────────────────────────────
# Slide 3: 表示モード切替
# ─────────────────────────────────────────────────────────
def slide_view_modes(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)

    # ヘッダー
    add_bg_rect(slide, 0, 0, 13.33, 1.1, C_DARK)
    add_bg_rect(slide, 0, 0, 0.18, 1.1, C_GOLD)
    add_text(slide, "表示モード切替ボタン", 0.4, 0.15, 10, 0.8,
             font_size=28, bold=True, color=C_WHITE)
    add_text(slide, "2 / 6", 12.0, 0.25, 1.0, 0.6,
             font_size=14, color=RGBColor(0xBA, 0xD4, 0xFF), align=PP_ALIGN.RIGHT)

    # ボタン説明カード 4枚
    modes = [
        ("👤 担当別",   "担当者別リソース全画面表示\n\n各担当者の工程を\n縦に並べて確認",
         C_DARK2,  C_WHITE),
        ("🔧 組立",     "組立タスクのみ表示\n（デフォルト表示）\n\n機械組立・盤組立・\n電気艤装などを確認",
         C_TEAL,   C_WHITE),
        ("📍 組立場所", "フロアプランビュー\n\n工場の区画ごとに\n作業場所を確認",
         C_BLUE,   C_WHITE),
        ("✈ 出張",      "出張タスクのみ表示\n\n出張スケジュールを\n専用画面で確認",
         RGBColor(0x55, 0x20, 0x80),  C_WHITE),
    ]

    card_x = [0.4, 3.7, 6.97, 10.2]
    cw, ch = 3.1, 4.7

    for i, (label, desc, bg, fg) in enumerate(modes):
        cx = card_x[i]
        add_card(slide, cx, 1.25, cw, ch, bg=C_WHITE)

        # カード上部色帯
        add_bg_rect(slide, cx, 1.25, cw, 1.05, bg)
        add_text(slide, label, cx, 1.30, cw, 0.95,
                 font_size=18, bold=True, color=fg,
                 align=PP_ALIGN.CENTER)

        # 説明
        add_text(slide, desc,
                 cx + 0.15, 2.45, cw - 0.3, 3.2,
                 font_size=13.5, color=C_GRAY,
                 align=PP_ALIGN.LEFT, wrap=True)

    # ヒント
    add_bg_rect(slide, 0.4, 6.2, 12.5, 0.9, C_BLUE)
    add_text(slide,
             "💡  ボタンをクリックするとモード切替。アクティブなボタンは強調表示されます。"
             "同じボタンを再クリックすると「組立」モードに戻ります。",
             0.6, 6.3, 12.1, 0.7,
             font_size=13, color=C_WHITE, wrap=True)

    return slide


# ─────────────────────────────────────────────────────────
# Slide 4: ズーム・フィルター
# ─────────────────────────────────────────────────────────
def slide_zoom_filter(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)

    add_bg_rect(slide, 0, 0, 13.33, 1.1, C_DARK)
    add_bg_rect(slide, 0, 0, 0.18, 1.1, C_GOLD)
    add_text(slide, "ズーム・フィルター操作", 0.4, 0.15, 10, 0.8,
             font_size=28, bold=True, color=C_WHITE)
    add_text(slide, "3 / 6", 12.0, 0.25, 1.0, 0.6,
             font_size=14, color=RGBColor(0xBA, 0xD4, 0xFF), align=PP_ALIGN.RIGHT)

    # ─ ズーム（左側）─
    add_card(slide, 0.4, 1.25, 5.8, 5.55, bg=C_WHITE)
    add_bg_rect(slide, 0.4, 1.25, 5.8, 0.55, C_TEAL)
    add_text(slide, "ズーム切替", 0.55, 1.28, 5.5, 0.48,
             font_size=16, bold=True, color=C_WHITE)

    items_zoom = [
        ("日単位", "1日1列でバーを詳細表示\n（デフォルト）\n土日・祝日は背景色が変わります"),
        ("週単位", "1週間を1列に圧縮\n全体スケジュールを\nコンパクトに把握"),
    ]
    for j, (btn, desc) in enumerate(items_zoom):
        bx = 0.7 + j * 2.7
        add_bg_rect(slide, bx, 2.0, 2.2, 0.5, C_DARK2)
        add_text(slide, btn, bx, 2.0, 2.2, 0.5,
                 font_size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc,
                 bx, 2.62, 2.3, 1.2,
                 font_size=12.5, color=C_GRAY, wrap=True)

    # 今日ライン
    add_card(slide, 0.55, 4.1, 5.5, 2.5, bg=RGBColor(0xF8, 0xFB, 0xFF),
             border_color=RGBColor(0xCC, 0xDD, 0xEE))
    add_text(slide, "赤い縦線 = 今日の位置",
             0.7, 4.2, 5.2, 0.5,
             font_size=13, bold=True, color=C_RED)
    add_text(slide,
             "タイムライン上の赤い縦線が「本日」を示します。\n"
             "現在進行中のタスクを素早く特定できます。",
             0.7, 4.75, 5.2, 1.5,
             font_size=12.5, color=C_GRAY, wrap=True)

    # ─ フィルター（右側）─
    add_card(slide, 6.5, 1.25, 6.5, 5.55, bg=C_WHITE)
    add_bg_rect(slide, 6.5, 1.25, 6.5, 0.55, C_BLUE)
    add_text(slide, "フィルター", 6.65, 1.28, 6.2, 0.48,
             font_size=16, bold=True, color=C_WHITE)

    filters = [
        ("工事番号フィルター",
         "「工事番号: 全表示」をクリックしてドロップダウンを開きます。\n"
         "工事番号を1つ以上選択して絞り込みが可能。複数選択もできます。"),
        ("タスク名フィルター",
         "「タスク名: 全表示」をクリックで絞り込み。\n"
         "機械組立 / 盤組立 / 電気艤装 / 出荷 の4種類から選択できます。"),
        ("担当者フィルター",
         "「リソース表示」をONにすると表示されます。\n"
         "特定の担当者のタスクのみ表示する絞り込みが可能です。"),
    ]

    fy = 1.95
    for ftitle, fdesc in filters:
        add_bg_rect(slide, 6.6, fy, 0.06, 1.35, C_GOLD)
        add_text(slide, ftitle,
                 6.8, fy, 6.0, 0.48,
                 font_size=13, bold=True, color=C_DARK)
        add_text(slide, fdesc,
                 6.8, fy + 0.48, 6.0, 0.9,
                 font_size=12, color=C_GRAY, wrap=True)
        fy += 1.55

    # フッター
    add_bg_rect(slide, 0.4, 7.0, 12.5, 0.38, RGBColor(0xE8, 0xF0, 0xF8))
    add_text(slide,
             "💡  フィルター外をクリックするとドロップダウンが閉じます",
             0.6, 7.05, 12.1, 0.3,
             font_size=12, color=C_TEAL)

    return slide


# ─────────────────────────────────────────────────────────
# Slide 5: グリッド操作
# ─────────────────────────────────────────────────────────
def slide_grid_ops(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)

    add_bg_rect(slide, 0, 0, 13.33, 1.1, C_DARK)
    add_bg_rect(slide, 0, 0, 0.18, 1.1, C_GOLD)
    add_text(slide, "グリッド（左側）の操作", 0.4, 0.15, 10, 0.8,
             font_size=28, bold=True, color=C_WHITE)
    add_text(slide, "4 / 6", 12.0, 0.25, 1.0, 0.6,
             font_size=14, color=RGBColor(0xBA, 0xD4, 0xFF), align=PP_ALIGN.RIGHT)

    ops = [
        (
            "シングルクリック",
            "タスクの開始日へ自動スクロール",
            "グリッドの行（タスク）をクリックすると、\n"
            "タイムライン（右側）がそのタスクの\n"
            "開始日の位置まで自動的にスクロールします。\n\n"
            "探したいタスクをすぐに見つけられます。",
            C_TEAL, "🖱",
        ),
        (
            "ダブルクリック",
            "タスク詳細の表示",
            "グリッドの行をダブルクリックすると、\n"
            "インライン編集ポップアップが開きます。\n\n"
            "閲覧者（ログアウト状態）は\n"
            "詳細内容の確認のみ可能です。\n"
            "変更はできません。",
            C_BLUE, "✏",
        ),
        (
            "Ctrl + クリック",
            "複数行の選択（参照用）",
            "Ctrlキーを押しながら複数行をクリックすると\n"
            "複数のタスクを同時に選択できます。\n\n"
            "閲覧者は選択のみ可能です。\n"
            "削除・編集操作は行えません。",
            C_DARK2, "☑",
        ),
    ]

    cx_list = [0.4, 4.7, 8.97]
    cw, ch = 4.0, 5.55

    for i, (op_name, op_sub, op_desc, op_color, icon) in enumerate(ops):
        cx = cx_list[i]
        add_card(slide, cx, 1.25, cw, ch, bg=C_WHITE)

        # カラー帯
        add_bg_rect(slide, cx, 1.25, cw, 0.95, op_color)

        # アイコン
        add_text(slide, icon, cx + 0.1, 1.3, 0.8, 0.85,
                 font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # タイトル
        add_text(slide, op_name,
                 cx + 0.9, 1.28, cw - 1.0, 0.52,
                 font_size=15, bold=True, color=C_WHITE)
        add_text(slide, op_sub,
                 cx + 0.9, 1.72, cw - 1.0, 0.42,
                 font_size=11, color=RGBColor(0xCC, 0xEE, 0xFF))

        # 説明
        add_text(slide, op_desc,
                 cx + 0.2, 2.3, cw - 0.4, 4.2,
                 font_size=13.5, color=C_GRAY, wrap=True)

    # フッター
    add_bg_rect(slide, 0.4, 6.95, 12.5, 0.42, C_BLUE)
    add_text(slide,
             "※  タスクバーのドラッグ（期間変更）は編集者のみ可能です。閲覧者はスクロール確認のみです。",
             0.6, 7.0, 12.1, 0.35,
             font_size=12, color=C_GOLD)

    return slide


# ─────────────────────────────────────────────────────────
# Slide 6: リソース表示
# ─────────────────────────────────────────────────────────
def slide_resource(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)

    add_bg_rect(slide, 0, 0, 13.33, 1.1, C_DARK)
    add_bg_rect(slide, 0, 0, 0.18, 1.1, C_GOLD)
    add_text(slide, "リソース表示 ／ 担当別モード", 0.4, 0.15, 10, 0.8,
             font_size=28, bold=True, color=C_WHITE)
    add_text(slide, "5 / 6", 12.0, 0.25, 1.0, 0.6,
             font_size=14, color=RGBColor(0xBA, 0xD4, 0xFF), align=PP_ALIGN.RIGHT)

    # 左側: リソース表示
    add_card(slide, 0.4, 1.25, 5.9, 5.6, bg=C_WHITE)
    add_bg_rect(slide, 0.4, 1.25, 5.9, 0.55, C_TEAL)
    add_text(slide, "🔲 リソース表示（ガント下部に追加）",
             0.55, 1.28, 5.6, 0.48,
             font_size=14, bold=True, color=C_WHITE)

    res_items = [
        ("「リソース表示」ボタンをクリック",
         "ガントチャートの下部に担当者別パネルが追加表示されます。"),
        ("担当者ごとの業務量を確認",
         "各担当者がどの期間にどのタスクを担当しているか\n一覧で確認できます。"),
        ("タイムラインと連動スクロール",
         "ガントチャートと同期して左右スクロールします。"),
    ]
    ry = 1.95
    for title, desc in res_items:
        add_bg_rect(slide, 0.6, ry + 0.12, 0.06, 0.9, C_TEAL)
        add_text(slide, title,
                 0.8, ry, 5.3, 0.48,
                 font_size=13, bold=True, color=C_DARK)
        add_text(slide, desc,
                 0.8, ry + 0.44, 5.3, 0.75,
                 font_size=12.5, color=C_GRAY, wrap=True)
        ry += 1.4

    # 右側: 担当別モード
    add_card(slide, 6.7, 1.25, 6.3, 5.6, bg=C_WHITE)
    add_bg_rect(slide, 6.7, 1.25, 6.3, 0.55, C_DARK2)
    add_text(slide, "👤 担当別（全画面リソース表示）",
             6.85, 1.28, 6.0, 0.48,
             font_size=14, bold=True, color=C_WHITE)

    res2_items = [
        ("「👤 担当別」ボタンをクリック",
         "ガントチャートが非表示になり、\n担当者別のリソース画面が全画面で開きます。"),
        ("担当者名をクリック",
         "その担当者のタスク一覧が\nタスクタイプ別に詳細表示されます。"),
        ("「← 一覧に戻る」ボタン",
         "担当者一覧に戻ります。"),
        ("再度ボタンをクリック",
         "通常のガントチャート（組立モード）に\n戻ります。"),
    ]
    ry2 = 1.95
    for title2, desc2 in res2_items:
        add_bg_rect(slide, 6.9, ry2 + 0.12, 0.06, 0.9, C_GOLD)
        add_text(slide, title2,
                 7.1, ry2, 5.7, 0.48,
                 font_size=13, bold=True, color=C_DARK)
        add_text(slide, desc2,
                 7.1, ry2 + 0.44, 5.7, 0.75,
                 font_size=12.5, color=C_GRAY, wrap=True)
        ry2 += 1.35

    return slide


# ─────────────────────────────────────────────────────────
# Slide 7: 使い方ガイド・ログイン
# ─────────────────────────────────────────────────────────
def slide_help_login(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg_rect(slide, 0, 0, 13.33, 7.5, C_LIGHT)

    add_bg_rect(slide, 0, 0, 13.33, 1.1, C_DARK)
    add_bg_rect(slide, 0, 0, 0.18, 1.1, C_GOLD)
    add_text(slide, "使い方ガイド ／ ログインについて", 0.4, 0.15, 10, 0.8,
             font_size=28, bold=True, color=C_WHITE)
    add_text(slide, "6 / 6", 12.0, 0.25, 1.0, 0.6,
             font_size=14, color=RGBColor(0xBA, 0xD4, 0xFF), align=PP_ALIGN.RIGHT)

    # 左: 使い方ガイド
    add_card(slide, 0.4, 1.25, 5.9, 5.55, bg=C_WHITE)
    add_bg_rect(slide, 0.4, 1.25, 5.9, 0.6, C_TEAL)
    add_text(slide, "？ 使い方ガイドボタン",
             0.55, 1.3, 5.6, 0.5,
             font_size=15, bold=True, color=C_WHITE)

    add_text(slide,
             "画面右上の「？使い方ガイド」ボタンを\n"
             "クリックすると、各ボタンの説明が\n"
             "ポップアップで表示されます。\n\n"
             "各ボタンにカーソルを重ねると\n"
             "詳細な説明が表示されます。\n\n"
             "もう一度クリックするか、\n"
             "Escキーで閉じることができます。",
             0.6, 1.98, 5.5, 4.0,
             font_size=14, color=C_GRAY, wrap=True)

    # 右: ログインについて
    add_card(slide, 6.7, 1.25, 6.3, 5.55, bg=C_WHITE)
    add_bg_rect(slide, 6.7, 1.25, 6.3, 0.6, C_DARK2)
    add_text(slide, "ログインについて",
             6.85, 1.3, 6.0, 0.5,
             font_size=15, bold=True, color=C_WHITE)

    # 閲覧者でできること
    add_bg_rect(slide, 6.85, 1.98, 6.0, 0.38, RGBColor(0xE8, 0xF5, 0xF0))
    add_text(slide, "✅  閲覧者（ログアウト状態）でできること",
             7.0, 2.0, 5.8, 0.34,
             font_size=12, bold=True, color=C_GREEN)

    viewer_can = [
        "工程表の閲覧・スクロール",
        "フィルターによる絞り込み",
        "表示モードの切替",
        "リソース表示の確認",
        "担当者別タスクの確認",
    ]
    vy = 2.43
    for item in viewer_can:
        add_text(slide, "  ✓  " + item,
                 7.0, vy, 5.8, 0.38,
                 font_size=12.5, color=C_GREEN)
        vy += 0.38

    # 要ログイン
    add_bg_rect(slide, 6.85, vy + 0.1, 6.0, 0.38, RGBColor(0xFE, 0xF2, 0xF2))
    add_text(slide, "🔒  編集者（ログイン後）のみできること",
             7.0, vy + 0.12, 5.8, 0.34,
             font_size=12, bold=True, color=C_RED)

    editor_only = [
        "タスクの追加・編集・削除",
        "ガントバーのドラッグ操作",
        "完了工事のアーカイブ",
    ]
    ey = vy + 0.55
    for item in editor_only:
        add_text(slide, "  🔒  " + item,
                 7.0, ey, 5.8, 0.38,
                 font_size=12.5, color=C_RED)
        ey += 0.38

    # ログイン手順
    add_bg_rect(slide, 6.85, ey + 0.1, 6.0, 0.9, RGBColor(0xE0, 0xEE, 0xFF))
    add_text(slide,
             "ログイン方法：右上「ログイン」→ メールアドレスとパスワードを入力",
             7.0, ey + 0.18, 5.8, 0.75,
             font_size=12, color=C_DARK, wrap=True)

    return slide


# ─────────────────────────────────────────────────────────
# メイン
# ─────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_title(prs)
    slide_overview(prs)
    slide_view_modes(prs)
    slide_zoom_filter(prs)
    slide_grid_ops(prs)
    slide_resource(prs)
    slide_help_login(prs)

    out = r"C:\Users\kurosaki\Desktop\工程表作成\組立工程表\組立工程表_閲覧者向け操作ガイド.pptx"
    prs.save(out)
    print(f"保存完了: {out}")


if __name__ == "__main__":
    main()
